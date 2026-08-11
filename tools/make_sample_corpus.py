#!/usr/bin/env python3
"""Generate a small synthetic data room so you can exercise the app end to end.

    python3 tools/make_sample_corpus.py ./sample-dataroom

Produces PDFs, a multi-sheet workbook with formulas, a legacy Excel 97-2003
workbook, a deck with speaker notes, a Word contract, a CSV, and a deliberate
near-duplicate pair — everything the ingest and dedupe paths need to be worth
looking at. All figures are invented.
"""

from __future__ import annotations

import base64
import gzip
import sys
from pathlib import Path

# A genuine Excel 97-2003 workbook (OLE2 container, BIFF8 records), embedded
# rather than generated: writing one needs xlwt, and a legacy-format *writer* is
# not worth a dependency for an app that only ever reads them. Two sheets,
# "Consolidated P&L" and "Notes", with a date-formatted cell and a boolean so the
# legacy extractor's rendering is exercised, not just its happy path.
#
# Regenerate with:  pip install xlwt  (then xlwt.Workbook(); ...; wb.save(buf);
# base64.b64encode(gzip.compress(buf.getvalue(), mtime=0)))
_LEGACY_XLS_GZ_B64 = (
    "H4sIAAAAAAACA+1YTWgUSRR+1dOdmWh+JmOUjUIoAsZ/ifEigo66+WHBZMeg+MOCO5OpGRsn3UNXjz972NVoji6C"
    "J5e9CNnDXtyfw+qiol7Eg6DoQVhYSHZv62lBQVDTvnrdLYm2YFCDSr+mXlW9el+9N1WvX9X0ndstE+d+WzwJL9Fm"
    "SMCUVw9102QMS33YSQOOe55qhnUKixfTR0X1KdzIOgMuNd5Kqj1U+z0JGvyqX0cO8A+Wr6AKg7Yl+BzSNvIhz5QP"
    "m5iOsafBj1iaoI38yhAfJr6A+C+ke5n4FpKcIr4JdSfYPridHVy5IYjjPVoHjTUhZ3CBMH+RZB0shJsqjr87zXxd"
    "A7Y6Zr7yAQwwHLisveVU7XoDjAPudb+whKP6zfATzAPYi7RmYGBNT88EtGIcjMMjjwM8DF/4azyWz62cAcofz5Qn"
    "XyNPvUZej5v7svyMpgMcA+9rpqJ/DD6DnOGn9M9tS9oVs5h3RZHnOrePQSNsUGkdo2fQdoV8Cm2sAbsNfrrv3TXE"
    "l3V1dWHW6Nvb3dXdHTbWYYANiUPCqolGwGmly+0Sl/mKkAjtd2wpedWxS6abAfiyioHomlaZiyNVYUkhcZrebV/s"
    "7NmK4O2iWBYOH67YUhRTalrpKv/6AHKOqOYddLVmFVHFPSBwUnHItGuSV3yYPCpdMcILomQ7gjTIv/V8xCwro7a1"
    "ti90npfMcs0Rkh8QqJsP9DHxSWrka0VTrUt+eNiuWa7kh5VaoWZWXF5y7JG189T5SPk0PSOfNlKWaUBehGZqt1Cu"
    "SeMSPv35/7sDhVx2P0mO0Rnqn7RL1S6BB8cVAsFNNOKfvwksKwmxivgozbqE2ouJt+KOYd2ZWxg0+k6Qzkka7UQ7"
    "64nuZZdNay/H9tiDHRfbx/7NrsD2+f7Jb1rP38+egw60XES8ek7Aaraa/XBW0Z/ZsGZBTv6beNsr+TmlpQPfveA6"
    "0QzPVOJBaiHu99TqsBc9LVgrhWYRaEZngVqPK1gr79ALewHyu6MKqXBaBM4/Q/QAp/m43Q+9DIyWH3k+LhGBSxDO"
    "CHAJH3cjiXzov7rAnh6B0wlXF+B0wnVsfIL2CuNPAntGBM4gXDLAGT4OtTLw+/daYC9J6zQTlySciqZvNdVjuAvt"
    "V339VIR+ivRVdBlaCnuMVnOzloE/6MTZMu0eOBdxrr0Y/1jjnEXEeYbqMJZZZCzPV7dvte4Qte4xxRRTTDHF9C6I"
    "BWetutfodM/3z/Jk8F3nGZap+DPJJ0tDYOPj4j/EXrCwduDorOJnERgsnIu9ISb8XqhoN1p34CAUyI+Ds45fvIux"
    "6b/njYHpd/cKzdb+1Gz8fM/2nwObFWk/ABYAAA=="
)


def make(root: Path) -> None:
    (root / "01_financial").mkdir(parents=True, exist_ok=True)
    (root / "02_legal").mkdir(parents=True, exist_ok=True)
    (root / "03_commercial").mkdir(parents=True, exist_ok=True)
    (root / "04_hr").mkdir(parents=True, exist_ok=True)

    # --- workbook with formulas ---
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "P&L"
    ws.append(["EUR '000", "FY2022", "FY2023", "FY2024"])
    ws.append(["Revenue", 318400, 372100, 412600])
    ws.append(["Cost of sales", -191040, -219539, -239308])
    ws.append(["Gross profit", "=C3+C4", "=D3+D4", "=E3+E4"])
    ws.append(["Operating expenses", -78200, -86400, -94100])
    ws.append(["EBITDA", "=C5+C6", "=D5+D6", "=E5+E6"])
    seg = wb.create_sheet("Segments")
    seg.append(["Segment", "FY2022", "FY2023", "FY2024"])
    for name, a, b, c in [
        ("Industrial", 141200, 163800, 178900),
        ("Consumer", 98700, 118300, 133400),
        ("Services", 78500, 90000, 100300),
    ]:
        seg.append([name, a, b, c])
    cust = wb.create_sheet("Customers")
    cust.append(["Customer", "FY2024 revenue", "Share of revenue"])
    for name, rev in [("Meridian Group", 128900), ("Halstead AG", 51200), ("Others", 232500)]:
        cust.append([name, rev, f"=B{cust.max_row + 1}/412600"])
    wb.save(root / "01_financial" / "Project_Kestrel_Financial_Model_v3.xlsx")

    # --- legacy workbook, the kind that has been sitting in the data room since
    # before anyone standardised on .xlsx ---
    (root / "01_financial" / "FinancialStatements_Consolidated_WorkingFile_2022_old.xls") \
        .write_bytes(gzip.decompress(base64.b64decode(_LEGACY_XLS_GZ_B64)))

    # --- PDFs ---
    import pymupdf

    def pdf(path: Path, pages: list[tuple[str, str]]) -> None:
        doc = pymupdf.open()
        for title, body in pages:
            page = doc.new_page()
            page.insert_text((60, 70), title, fontsize=15)
            y = 105
            for line in body.split("\n"):
                page.insert_text((60, y), line[:105], fontsize=9.5)
                y += 14
        doc.save(str(path))
        doc.close()

    pdf(
        root / "01_financial" / "Audited_Accounts_FY2024.pdf",
        [
            ("Kestrel Holding GmbH — Audited Financial Statements FY2024",
             "Independent auditor's report\n\n"
             "We have audited the financial statements of Kestrel Holding GmbH for the\n"
             "year ended 31 December 2024. In our opinion the financial statements give a\n"
             "true and fair view of the financial position of the Company.\n\n"
             "Basis of opinion: audit conducted in accordance with Austrian GAAP.\n"
             "Auditor: Brandmayr & Partner Wirtschaftsprüfung GmbH, Vienna."),
            ("Statement of profit or loss",
             "EUR '000                        FY2024      FY2023\n"
             "Revenue                          412,600     372,100\n"
             "Cost of sales                   (239,308)   (219,539)\n"
             "Gross profit                     173,292     152,561\n"
             "Operating expenses               (94,100)    (86,400)\n"
             "EBITDA                            79,192      66,161\n"
             "Depreciation and amortisation    (18,400)    (17,100)\n"
             "Operating profit                  60,792      49,061"),
            ("Notes to the financial statements — note 14, contingencies",
             "The Company is party to a tax audit covering financial years 2020 to 2022.\n"
             "Management estimates a reasonably possible additional assessment of up to\n"
             "EUR 4.2 million, for which no provision has been recognised.\n\n"
             "Note 15 — customer concentration\n"
             "One customer accounted for 31.2% of revenue in FY2024 (FY2023: 24.8%)."),
        ],
    )
    pdf(
        root / "02_legal" / "Supply_Agreement_Meridian_2021.pdf",
        [
            ("Master Supply Agreement — Kestrel Holding GmbH and Meridian Group plc",
             "Dated 14 March 2021.\n\n"
             "Clause 3.1 Term. This Agreement commences on the Effective Date and\n"
             "continues for five (5) years unless terminated earlier.\n\n"
             "Clause 11.2 Change of control. Either party may terminate this Agreement\n"
             "on thirty (30) days' written notice if the other party undergoes a change\n"
             "of control, being the acquisition of more than 50% of its voting shares.\n\n"
             "Clause 12.4 Limitation of liability. Aggregate liability shall not exceed\n"
             "the fees paid in the twelve months preceding the claim."),
            ("Schedule 2 — pricing and volume commitments",
             "Minimum annual volume: 240,000 units.\n"
             "Price escalation capped at CPI + 1.5% per annum.\n"
             "Rebate: 2.5% of annual spend above EUR 100 million.\n\n"
             "Signed for and on behalf of the parties. [SIGNATURE PAGE NOT INCLUDED]"),
        ],
    )
    pdf(
        root / "04_hr" / "Employment_Agreement_CEO_unsigned.pdf",
        [
            ("Service Agreement — Managing Director",
             "DRAFT — not for execution.\n\n"
             "Clause 4. Remuneration. Base salary of EUR 385,000 per annum.\n"
             "Clause 5. Bonus. Up to 60% of base salary on achievement of targets.\n"
             "Clause 9. Change of control payment. On a change of control the Executive\n"
             "is entitled to a payment equal to 18 months' base salary.\n"
             "Clause 12. Non-compete. Twelve months post-termination, EU-wide."),
        ],
    )

    # --- deck with notes ---
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Project Kestrel — Management Presentation"
    s.placeholders[1].text = "Confidential · March 2025"
    for title, bullets, notes in [
        ("Financial track record",
         ["Revenue EUR 412.6m in FY2024, +10.9% year on year",
          "EBITDA EUR 79.2m, 19.2% margin",
          "Three consecutive years of margin expansion"],
         "Note: FY2024 EBITDA shown here is pre-exceptional. The audited figure is the same "
          "79.2m but management add back 2.1m of restructuring in the plan case."),
        ("Growth plan 2025-2027",
         ["Revenue target EUR 560m by FY2027",
          "EBITDA margin target 22%",
          "Two bolt-on acquisitions assumed, not yet signed"],
         "The 560m target assumes the Meridian contract renews on current terms. No LOI exists "
          "for either bolt-on."),
        ("Customer base",
         ["Top customer 31% of FY2024 revenue",
          "Average tenure 7 years",
          "Contracted backlog EUR 214m"],
         "Do not volunteer the change-of-control clause in the Meridian agreement unless asked."),
    ]:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = title
        tf = sl.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.font.size = Pt(18)
        sl.notes_slide.notes_text_frame.text = notes
    prs.save(root / "03_commercial" / "Management_Presentation_March2025.pptx")

    # --- Word doc ---
    import docx

    d = docx.Document()
    d.add_heading("Vendor Due Diligence Report — Executive Summary", 0)
    d.add_paragraph(
        "This report has been prepared by the vendor's advisers. Findings are based on "
        "information provided by management and have not been independently verified."
    )
    d.add_heading("Key findings", level=1)
    for t in [
        "Revenue grew from EUR 318.4m in FY2022 to EUR 412.6m in FY2024.",
        "Customer concentration is elevated: the largest customer represents 31.2% of FY2024 revenue.",
        "An open tax audit covering FY2020-FY2022 is unprovided, with exposure up to EUR 4.2m.",
        "The largest supply agreement contains a change-of-control termination right.",
    ]:
        d.add_paragraph(t, style="List Bullet")
    d.save(str(root / "01_financial" / "VDD_Report_Executive_Summary.docx"))

    # --- CSV + a near-duplicate of the model ---
    (root / "03_commercial" / "monthly_bookings.csv").write_text(
        "month,bookings_eur_000,units\n"
        + "\n".join(
            f"2024-{m:02d},{28000 + m * 640},{19000 + m * 310}" for m in range(1, 13)
        )
        + "\n",
        encoding="utf-8",
    )
    # A near-duplicate of the model: same document, one figure revised. This is
    # the case that actually confuses a DD assistant, so it is worth having in
    # the fixture.
    ws["E6"] = -96_400
    wb.save(root / "01_financial" / "Project_Kestrel_Financial_Model_v3_FINAL.xlsx")
    # …and an exact copy filed in a second folder, which should collapse to one document.
    dup = (root / "01_financial" / "Project_Kestrel_Financial_Model_v3.xlsx").read_bytes()
    (root / "03_commercial").mkdir(exist_ok=True)
    (root / "03_commercial" / "Model_copy_from_financial.xlsx").write_bytes(dup)

    print(f"Sample data room written to {root.resolve()}")
    for p in sorted(root.rglob("*")):
        if p.is_file():
            print(f"  {p.relative_to(root)}  ({p.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    make(Path(sys.argv[1] if len(sys.argv) > 1 else "./sample-dataroom"))
