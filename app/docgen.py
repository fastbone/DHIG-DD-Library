"""Deliverable generation: DOCX / XLSX / PPTX / Markdown from a block spec.

The agent hands over a declarative spec rather than writing file-format code,
which keeps output deterministic and keeps citations intact.
"""

from __future__ import annotations

import json
import re
import time
import uuid
from pathlib import Path

from . import db
from .config import settings

SAFE_NAME = re.compile(r"[^A-Za-z0-9._-]+")

BLOCK_SCHEMA = {
    "type": "object",
    "properties": {
        "type": {"type": "string", "enum": ["heading", "paragraph", "bullets", "table", "pagebreak"]},
        "level": {"type": "integer"},
        "text": {"type": "string"},
        "items": {"type": "array", "items": {"type": "string"}},
        "header": {"type": "array", "items": {"type": "string"}},
        "rows": {"type": "array", "items": {"type": "array", "items": {"type": "string"}}},
    },
    "required": ["type"],
}


def safe_filename(name: str, default_ext: str) -> str:
    name = SAFE_NAME.sub("_", Path(name).name).strip("._") or "document"
    if not Path(name).suffix:
        name = f"{name}.{default_ext}"
    return name[:120]


def _register(kind: str, path: Path, qa_id: str | None) -> dict:
    art_id = uuid.uuid4().hex[:12]
    db.execute(
        "INSERT INTO artifacts(id, kind, filename, path, size_bytes, qa_id, created_at)"
        " VALUES(?,?,?,?,?,?,?)",
        (art_id, kind, path.name, str(path), path.stat().st_size, qa_id, time.time()),
    )
    return {
        "artifact_id": art_id,
        "kind": kind,
        "filename": path.name,
        "size_bytes": path.stat().st_size,
        "download_url": f"/api/artifacts/{art_id}/download",
    }


def _md(title: str, blocks: list[dict]) -> str:
    out = [f"# {title}", ""] if title else []
    for b in blocks:
        t = b.get("type")
        if t == "heading":
            out += ["", "#" * max(1, min(int(b.get("level", 2)), 6)) + " " + b.get("text", ""), ""]
        elif t == "paragraph":
            out += [b.get("text", ""), ""]
        elif t == "bullets":
            out += [f"- {i}" for i in b.get("items", [])] + [""]
        elif t == "table":
            header = b.get("header", [])
            rows = b.get("rows", [])
            if header:
                out += ["| " + " | ".join(header) + " |",
                        "| " + " | ".join(["---"] * len(header)) + " |"]
            for r in rows:
                out.append("| " + " | ".join(str(c) for c in r) + " |")
            out.append("")
        elif t == "pagebreak":
            out += ["", "---", ""]
    return "\n".join(out)


def write_docx(title: str, blocks: list[dict], path: Path) -> None:
    import docx
    from docx.enum.text import WD_BREAK

    d = docx.Document()
    if title:
        d.add_heading(title, level=0)
    for b in blocks:
        t = b.get("type")
        if t == "heading":
            d.add_heading(b.get("text", ""), level=max(1, min(int(b.get("level", 2)), 5)))
        elif t == "paragraph":
            d.add_paragraph(b.get("text", ""))
        elif t == "bullets":
            for item in b.get("items", []):
                d.add_paragraph(str(item), style="List Bullet")
        elif t == "table":
            header = b.get("header", [])
            rows = b.get("rows", [])
            ncols = len(header) or (len(rows[0]) if rows else 1)
            table = d.add_table(rows=1 if header else 0, cols=ncols)
            table.style = "Light Grid Accent 1"
            if header:
                for i, h in enumerate(header[:ncols]):
                    table.rows[0].cells[i].text = str(h)
            for r in rows:
                cells = table.add_row().cells
                for i, c in enumerate(list(r)[:ncols]):
                    cells[i].text = str(c)
        elif t == "pagebreak":
            d.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    d.save(str(path))


def write_xlsx(title: str, sheets: list[dict], path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    wb.remove(wb.active)
    for s in sheets or [{"name": "Sheet1", "header": [], "rows": []}]:
        name = SAFE_NAME.sub(" ", str(s.get("name", "Sheet")))[:31] or "Sheet"
        ws = wb.create_sheet(title=name)
        header = s.get("header") or []
        if header:
            ws.append([str(h) for h in header])
            for c in ws[1]:
                c.font = Font(bold=True)
            ws.freeze_panes = "A2"
        for row in s.get("rows") or []:
            ws.append(["" if c is None else c for c in row])
        widths: dict[int, int] = {}
        for row in ws.iter_rows(values_only=True):
            for i, v in enumerate(row, start=1):
                widths[i] = min(max(widths.get(i, 10), len(str(v or "")) + 2), 60)
        for i, w in widths.items():
            ws.column_dimensions[get_column_letter(i)].width = w
    if not wb.sheetnames:
        wb.create_sheet("Sheet1")
    wb.save(str(path))


def write_pptx(title: str, slides: list[dict], path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    if title:
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = title
        if len(s.placeholders) > 1:
            s.placeholders[1].text = time.strftime("%d %B %Y")
    for spec in slides or []:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = str(spec.get("title", ""))
        body = s.placeholders[1].text_frame
        body.clear()
        bullets = spec.get("bullets") or []
        for i, b in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(b)
            para.font.size = Pt(18)
        notes = spec.get("notes")
        if notes:
            s.notes_slide.notes_text_frame.text = str(notes)
    prs.save(str(path))


def create(spec: dict, qa_id: str | None = None) -> dict:
    kind = (spec.get("kind") or "md").lower()
    title = spec.get("title") or ""
    blocks = spec.get("blocks") or []
    filename = safe_filename(spec.get("filename") or f"{title or 'document'}", kind)
    stamp = time.strftime("%Y%m%d-%H%M%S")
    path = settings.artifacts_dir / f"{stamp}_{filename}"

    if kind == "docx":
        write_docx(title, blocks, path)
    elif kind == "xlsx":
        sheets = spec.get("sheets") or [
            {"name": "Sheet1",
             "header": next((b.get("header") for b in blocks if b.get("type") == "table"), []),
             "rows": next((b.get("rows") for b in blocks if b.get("type") == "table"), [])}
        ]
        write_xlsx(title, sheets, path)
    elif kind == "pptx":
        slides = spec.get("slides") or []
        if not slides and blocks:
            slides = [{"title": b.get("text", ""), "bullets": b.get("items", [])} for b in blocks]
        write_pptx(title, slides, path)
    elif kind in {"md", "markdown", "txt"}:
        path = path.with_suffix(".md")
        path.write_text(_md(title, blocks), encoding="utf-8")
        kind = "md"
    elif kind == "json":
        path.write_text(json.dumps(spec.get("data", spec), indent=2), encoding="utf-8")
    else:
        raise ValueError(f"unsupported kind: {kind}")
    return _register(kind, path, qa_id)
