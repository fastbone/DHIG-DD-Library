"""Turn binary office files into a text mirror with citation anchors.

Every extractor returns ``(markdown_text, units)`` where each unit is a
citable atom — a PDF page, a slide, a spreadsheet sheet — carrying an anchor
string (``p12``, ``slide4``, ``Sheet1!A1:H240``) and a char range into the
markdown so the reader can jump straight to it.

The markdown mirror is written to ``data/derived/<doc_id>.md``, which also
makes the whole corpus greppable with ordinary tools.
"""

from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path

from .config import settings

MAX_CELL = 200
MAX_SHEET_ROWS = 5000
MAX_SHEET_COLS = 80


class ExtractionError(RuntimeError):
    pass


@dataclass
class Unit:
    ordinal: int
    anchor: str
    kind: str
    text: str
    char_start: int = 0
    char_end: int = 0


class Builder:
    """Accumulates markdown while recording char offsets per unit."""

    def __init__(self) -> None:
        self.buf = io.StringIO()
        self.units: list[Unit] = []
        self._pos = 0

    def write(self, text: str) -> None:
        self.buf.write(text)
        self._pos += len(text)

    def unit(self, anchor: str, kind: str, body: str) -> None:
        body = body.strip()
        if not body:
            return
        header = f"\n<!-- anchor: {anchor} -->\n"
        self.write(header)
        start = self._pos
        self.write(body + "\n")
        self.units.append(
            Unit(
                ordinal=len(self.units),
                anchor=anchor,
                kind=kind,
                text=body,
                char_start=start,
                char_end=self._pos,
            )
        )

    def finish(self) -> tuple[str, list[Unit]]:
        return self.buf.getvalue(), self.units


# --- PDF -----------------------------------------------------------------


def extract_pdf(path: Path, ocr: bool = False) -> tuple[str, list[Unit], bool]:
    import pymupdf

    b = Builder()
    ocr_used = False
    with pymupdf.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            if ocr and len(text.strip()) < 20:
                try:
                    tp = page.get_textpage_ocr(flags=0, full=True)
                    text = page.get_text("text", textpage=tp) or text
                    ocr_used = True
                except Exception:
                    pass
            b.unit(f"p{i}", "page", text)
    return (*b.finish(), ocr_used)


# --- PPTX ----------------------------------------------------------------


def extract_pptx(path: Path) -> tuple[str, list[Unit]]:
    from pptx import Presentation

    prs = Presentation(str(path))
    b = Builder()
    for i, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    chunks.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    chunks.append(" | ".join(cells))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            chunks.append(f"[speaker notes] {notes}")
        b.unit(f"slide{i}", "slide", "\n".join(chunks))
    return b.finish()


# --- Spreadsheets --------------------------------------------------------


def _cell(v) -> str:
    if v is None:
        return ""
    s = str(v)
    return s if len(s) <= MAX_CELL else s[:MAX_CELL] + "…"


def extract_xlsx(path: Path) -> tuple[str, list[Unit]]:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    b = Builder()
    # Opened by handle rather than by name, deliberately. openpyxl validates the
    # *file extension* before looking at the bytes, so a genuinely modern
    # workbook that someone saved as .xls is refused with a message about the old
    # format — and data rooms are full of exactly that. A file object skips the
    # name check; the container has already been identified by its magic bytes.
    #
    # Two passes: computed values (what an analyst reads) and formulas (how the
    # model was built). Both matter in a financial DD.
    with path.open("rb") as fh_val:
        wb_val = load_workbook(fh_val, read_only=True, data_only=True)
        try:
            formulas: dict[str, list[str]] = {}
            try:
                with path.open("rb") as fh_f:
                    wb_f = load_workbook(fh_f, read_only=True, data_only=False)
                    try:
                        for ws in wb_f.worksheets:
                            found: list[str] = []
                            for row in ws.iter_rows(
                                max_row=MAX_SHEET_ROWS, max_col=MAX_SHEET_COLS
                            ):
                                for c in row:
                                    if isinstance(c.value, str) and c.value.startswith("="):
                                        found.append(f"{c.coordinate}: {c.value[:MAX_CELL]}")
                                        if len(found) >= 300:
                                            break
                                if len(found) >= 300:
                                    break
                            if found:
                                formulas[ws.title] = found
                    finally:
                        wb_f.close()
            except Exception:  # noqa: BLE001 — values still index without formulas
                pass

            _xlsx_sheets(b, wb_val, formulas, get_column_letter)
        finally:
            wb_val.close()
    return b.finish()


def _xlsx_sheets(b: Builder, wb_val, formulas: dict[str, list[str]], get_column_letter) -> None:
    """The value pass: one citable unit per sheet, formulas appended where found."""
    for ws in wb_val.worksheets:
        lines: list[str] = []
        max_c = 0
        n_rows = 0
        for r_i, row in enumerate(
            ws.iter_rows(max_row=MAX_SHEET_ROWS, max_col=MAX_SHEET_COLS, values_only=True),
            start=1,
        ):
            cells = [_cell(v) for v in row]
            while cells and cells[-1] == "":
                cells.pop()
            if not cells:
                continue
            max_c = max(max_c, len(cells))
            n_rows = r_i
            lines.append(f"r{r_i}: " + " | ".join(cells))
        if not lines and ws.title not in formulas:
            continue
        span = f"A1:{get_column_letter(max(max_c, 1))}{max(n_rows, 1)}"
        body = [f"# sheet: {ws.title}  (range {span})", *lines]
        if ws.title in formulas:
            body.append("")
            body.append("## formulas")
            body.extend(formulas[ws.title])
        b.unit(f"{ws.title}!{span}", "sheet", "\n".join(body))


# --- legacy XLS (BIFF, Excel 97-2003) ------------------------------------


def _column_letter(index: int) -> str:
    """1 -> A. openpyxl's helper, without importing openpyxl for a legacy file."""
    letters = ""
    while index > 0:
        index, rem = divmod(index - 1, 26)
        letters = chr(65 + rem) + letters
    return letters or "A"


def extract_xls(path: Path) -> tuple[str, list[Unit]]:
    """Excel 97-2003 workbooks, via xlrd.

    openpyxl refuses these outright, and a data room that has been in use for a
    decade is full of them. Same anchor shape as the modern extractor
    (``Sheet1!A1:H240``) so a citation into a legacy file behaves identically.

    Two differences are inherent to the format and worth being explicit about in
    the text rather than silently absent:

    * xlrd exposes computed values only — the formula behind a cell is not
      recoverable — so the "## formulas" section the modern extractor writes has
      no equivalent here, and the mirror says so. Otherwise the model would read
      its absence as "this workbook has no formulas", which is a different and
      wrong conclusion.
    * dates are stored as serial numbers. Rendered as ISO dates here, because
      "44999" in a mirror is exactly the kind of figure that gets copied into a
      report as if it meant something.
    """
    import xlrd

    try:
        # logfile= keeps xlrd's running commentary out of the app log. It prints
        # notes like "*** No CODEPAGE record ... will use 'iso-8859-1'" straight to
        # stdout for perfectly readable files, which in a 2000-document sweep is
        # thousands of lines saying nothing actionable.
        book = xlrd.open_workbook(
            str(path), formatting_info=False, on_demand=True, logfile=io.StringIO()
        )
    except Exception as exc:  # noqa: BLE001 — xlrd raises several unrelated types
        raise ExtractionError(f"not a readable Excel 97-2003 workbook: {exc}") from exc

    b = Builder()
    try:
        for name in book.sheet_names():
            try:
                sh = book.sheet_by_name(name)
            except Exception:  # noqa: BLE001 — one bad sheet must not lose the rest
                continue
            lines: list[str] = []
            max_c = 0
            n_rows = 0
            for r_i in range(min(sh.nrows, MAX_SHEET_ROWS)):
                cells: list[str] = []
                for c_i in range(min(sh.ncols, MAX_SHEET_COLS)):
                    cells.append(_xls_cell(sh, r_i, c_i, book.datemode))
                while cells and cells[-1] == "":
                    cells.pop()
                if not cells:
                    continue
                max_c = max(max_c, len(cells))
                n_rows = r_i + 1
                lines.append(f"r{r_i + 1}: " + " | ".join(cells))
            if not lines:
                continue
            span = f"A1:{_column_letter(max(max_c, 1))}{max(n_rows, 1)}"
            body = [
                f"# sheet: {name}  (range {span})",
                *lines,
                "",
                "<!-- legacy .xls: cell formulas are not recoverable from this format;"
                " the values above are the stored computed values. -->",
            ]
            b.unit(f"{name}!{span}", "sheet", "\n".join(body))
            if book.on_demand:
                book.unload_sheet(name)
    finally:
        book.release_resources()
    return b.finish()


def _xls_cell(sh, row: int, col: int, datemode: int) -> str:
    """One legacy cell as text, with dates and errors rendered rather than raw."""
    import xlrd

    kind = sh.cell_type(row, col)
    value = sh.cell_value(row, col)
    if kind in (xlrd.XL_CELL_EMPTY, xlrd.XL_CELL_BLANK):
        return ""
    if kind == xlrd.XL_CELL_DATE:
        try:
            stamp = xlrd.xldate_as_datetime(value, datemode)
        except (ValueError, OverflowError, xlrd.XLDateError):
            return _cell(value)
        # A serial below 1 is a time of day with no date part — an `hh:mm` cell.
        # xldate_as_datetime lands those on the epoch day, so rendering the whole
        # timestamp turns a shift start of 15:00 into "1899-12-31 15:00:00", which
        # is precisely the invented figure this branch exists to prevent.
        if 0 <= float(value) < 1:
            return stamp.strftime("%H:%M:%S" if stamp.second else "%H:%M")
        # A serial with no time component is a plain date; keep it short.
        if (stamp.hour, stamp.minute, stamp.second) == (0, 0, 0):
            return stamp.date().isoformat()
        return stamp.isoformat(sep=" ")
    if kind == xlrd.XL_CELL_BOOLEAN:
        return "TRUE" if value else "FALSE"
    if kind == xlrd.XL_CELL_ERROR:
        return xlrd.error_text_from_code.get(value, f"#ERR{value}")
    if kind == xlrd.XL_CELL_NUMBER and float(value).is_integer():
        # 372100000.0 reads as a corrupted figure; 372100000 reads as the figure.
        return _cell(int(value))
    return _cell(value)


# --- DOCX ----------------------------------------------------------------


def extract_docx(path: Path) -> tuple[str, list[Unit]]:
    import docx

    d = docx.Document(str(path))
    b = Builder()
    chunk: list[str] = []
    section = 1
    CHARS = 6000
    size = 0
    for p in d.paragraphs:
        t = p.text.strip()
        if not t:
            continue
        chunk.append(t)
        size += len(t)
        if size >= CHARS:
            b.unit(f"sec{section}", "section", "\n".join(chunk))
            section += 1
            chunk, size = [], 0
    for ti, table in enumerate(d.tables, start=1):
        rows_txt = [
            " | ".join(c.text.strip().replace("\n", " ") for c in row.cells) for row in table.rows
        ]
        chunk.append(f"[table {ti}]\n" + "\n".join(rows_txt))
        size += sum(len(r) for r in rows_txt)
        if size >= CHARS:
            b.unit(f"sec{section}", "section", "\n".join(chunk))
            section += 1
            chunk, size = [], 0
    if chunk:
        b.unit(f"sec{section}", "section", "\n".join(chunk))
    return b.finish()


# --- Plain text / CSV ----------------------------------------------------


def extract_text(path: Path) -> tuple[str, list[Unit]]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    b = Builder()
    if path.suffix.lower() in {".csv", ".tsv"}:
        delim = "\t" if path.suffix.lower() == ".tsv" else ","
        reader = csv.reader(io.StringIO(raw), delimiter=delim)
        lines = [" | ".join(_cell(v) for v in row) for row in reader]
        for i in range(0, max(len(lines), 1), 500):
            b.unit(f"rows{i + 1}-{i + 500}", "rows", "\n".join(lines[i : i + 500]))
    elif path.suffix.lower() == ".json":
        try:
            pretty = json.dumps(json.loads(raw), indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            pretty = raw
        b.unit("json", "text", pretty[:400_000])
    else:
        step = 8000
        for i in range(0, max(len(raw), 1), step):
            b.unit(f"chars{i}-{i + step}", "text", raw[i : i + step])
    return b.finish()


# --- dispatch ------------------------------------------------------------


_OOXML_MAGIC = b"PK\x03\x04"          # .xlsx/.docx/.pptx are zip containers
_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"   # .xls/.doc/.ppt (BIFF/CFB)


def _container(path: Path) -> str:
    """What the file actually is: "ooxml", "ole2", or "other"."""
    try:
        with path.open("rb") as fh:
            head = fh.read(8)
    except OSError:
        return "other"
    if head.startswith(_OOXML_MAGIC):
        return "ooxml"
    if head.startswith(_OLE2_MAGIC):
        return "ole2"
    return "other"


def extract_spreadsheet(path: Path) -> tuple[str, list[Unit]]:
    """Route a spreadsheet by what it is rather than by what it is called.

    Extensions in a data room are unreliable in both directions: `.xls` files
    that are really modern workbooks, `.xlsx` files that are really BIFF, and —
    very commonly — "Excel" files that are actually HTML tables or CSV emitted by
    a reporting system. Dispatching on the extension turns each of those into a
    failed document; dispatching on the container gets all three indexed.
    """
    kind = _container(path)
    if kind == "ooxml":
        return extract_xlsx(path)
    if kind == "ole2":
        return extract_xls(path)
    # Neither container. Two things land here, in this order of likelihood:
    #
    # An HTML or delimited export named .xls — what reporting systems emit and
    # call a spreadsheet. Indexed as text, because the figures in it are what
    # someone will ask about.
    #
    # Or a pre-1995 workbook: BIFF2-4 is a bare record stream with no OLE2
    # wrapper, so the magic bytes above do not see it. xlrd still reads those, so
    # try it first and only fall back to text.
    try:
        return extract_xls(path)
    except ExtractionError:
        return extract_text(path)


def extract(path: Path, family: str, ocr: bool = False) -> tuple[str, list[Unit], bool]:
    if family == "pdf":
        return extract_pdf(path, ocr=ocr)
    if family == "pptx":
        return (*extract_pptx(path), False)
    if family == "xlsx":
        return (*extract_spreadsheet(path), False)
    if family == "docx":
        return (*extract_docx(path), False)
    if family == "text":
        return (*extract_text(path), False)
    raise ExtractionError(f"unsupported family: {family}")


def write_mirror(doc_id: str, text: str) -> Path:
    p = settings.text_path(doc_id)
    p.write_text(text, encoding="utf-8")
    return p


def read_mirror(doc_id: str, start: int = 0, end: int | None = None) -> str:
    p = settings.text_path(doc_id)
    if not p.exists():
        return ""
    text = p.read_text(encoding="utf-8", errors="replace")
    return text[start:end] if (start or end) else text
